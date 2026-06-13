from pptx import Presentation
import PyPDF2

def extract_ppt_text(file_path: str) -> list:
    """提取PPT每页的文本，返回列表，每项为一页内容"""
    prs = Presentation(file_path)
    slides = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    texts.append(para.text)
        slides.append('\n'.join(texts))
    return slides

def extract_pdf_text(file_path: str) -> list:
    """提取PDF每页文本，返回列表，每项为一页内容"""
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        return [page.extract_text() for page in reader.pages]